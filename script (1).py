# Continue with the remaining files and create the complete zip
all_files_continued = {
    "backend/services/ppt_generator.py": '''from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import json
from PIL import Image

class PPTGenerator:
    def __init__(self):
        self.themes = {
            "Corporate Blue": {
                "primary_color": RGBColor(44, 106, 160),
                "secondary_color": RGBColor(255, 255, 255),
                "accent_color": RGBColor(76, 175, 80),
                "text_color": RGBColor(33, 33, 33)
            },
            "Modern Green": {
                "primary_color": RGBColor(76, 175, 80),
                "secondary_color": RGBColor(255, 255, 255),
                "accent_color": RGBColor(255, 193, 7),
                "text_color": RGBColor(33, 33, 33)
            },
            "Creative Orange": {
                "primary_color": RGBColor(255, 152, 0),
                "secondary_color": RGBColor(255, 255, 255),
                "accent_color": RGBColor(233, 30, 99),
                "text_color": RGBColor(33, 33, 33)
            },
            "Minimal Gray": {
                "primary_color": RGBColor(96, 125, 139),
                "secondary_color": RGBColor(255, 255, 255),
                "accent_color": RGBColor(158, 158, 158),
                "text_color": RGBColor(33, 33, 33)
            }
        }
    
    def create_presentation(self, narrative, analysis_results, session_id, template_path=None, theme="Corporate Blue"):
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
            print(f"Using uploaded template: {template_path}")
        else:
            prs = Presentation()
            print(f"Using theme: {theme}")
        
        theme_colors = self.themes.get(theme, self.themes["Corporate Blue"])
        
        if template_path:
            slide_ids = [slide._element for slide in prs.slides][1:]
            for slide_id in slide_ids:
                prs.slides._sldIdLst.remove(slide_id)
        
        slides_data = narrative.get("slides", [])
        
        for i, slide_data in enumerate(slides_data):
            slide = self._create_slide(prs, slide_data, theme_colors, i == 0)
            self._add_charts_to_slide(slide, slide_data.get("chart_references", []))
        
        output_path = f"../data/outputs/{session_id}_presentation.pptx"
        prs.save(output_path)
        
        print(f"Presentation saved to: {output_path}")
        return output_path
    
    def _create_slide(self, prs, slide_data, theme_colors, is_title_slide=False):
        if is_title_slide:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = slide_data.get("title", "Data Analysis Report")
            self._format_title(title, theme_colors, size=Pt(44))
            
            if slide.shapes.placeholders[1] and slide_data.get("content"):
                subtitle = slide.shapes.placeholders[1]
                subtitle.text = "\\n".join(slide_data["content"][:2])
                self._format_text(subtitle, theme_colors, size=Pt(24))
        
        else:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = slide_data.get("title", "Analysis")
            self._format_title(title, theme_colors, size=Pt(32))
            
            content_placeholder = slide.shapes.placeholders[1]
            tf = content_placeholder.text_frame
            
            for i, bullet_point in enumerate(slide_data.get("content", [])):
                if i == 0:
                    tf.text = bullet_point
                else:
                    p = tf.add_paragraph()
                    p.text = bullet_point
                
                if i < len(tf.paragraphs):
                    self._format_paragraph(tf.paragraphs[i], theme_colors)
        
        return slide
    
    def _add_charts_to_slide(self, slide, chart_references):
        if not chart_references:
            return
        
        chart_left = Inches(6)
        chart_top = Inches(2)
        chart_width = Inches(4)
        chart_height = Inches(3)
        
        for i, chart_path in enumerate(chart_references[:2]):
            if os.path.exists(chart_path):
                try:
                    top_position = chart_top + (i * Inches(3.5))
                    slide.shapes.add_picture(
                        chart_path, 
                        chart_left, 
                        top_position, 
                        chart_width, 
                        chart_height
                    )
                    print(f"Added chart: {chart_path}")
                except Exception as e:
                    print(f"Error adding chart {chart_path}: {e}")
    
    def _format_title(self, title_shape, theme_colors, size=Pt(32)):
        if hasattr(title_shape, 'text_frame'):
            tf = title_shape.text_frame
            if tf.paragraphs:
                p = tf.paragraphs[0]
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = size
                run.font.color.rgb = theme_colors["primary_color"]
                run.font.bold = True
                p.alignment = PP_ALIGN.CENTER
    
    def _format_text(self, text_shape, theme_colors, size=Pt(18)):
        if hasattr(text_shape, 'text_frame'):
            tf = text_shape.text_frame
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.size = size
                    run.font.color.rgb = theme_colors["text_color"]
    
    def _format_paragraph(self, paragraph, theme_colors, size=Pt(18)):
        for run in paragraph.runs:
            run.font.size = size
            run.font.color.rgb = theme_colors["text_color"]
        
        paragraph.level = 0''',

    "backend/utils/file_utils.py": '''import aiofiles
import os
from fastapi import UploadFile
import pandas as pd

async def save_upload_file(upload_file: UploadFile, destination: str) -> str:
    """Save uploaded file to destination"""
    os.makedirs(os.path.dirname(destination), exist_ok=True)
    
    async with aiofiles.open(destination, 'wb') as f:
        content = await upload_file.read()
        await f.write(content)
    
    return destination

def validate_csv_file(file: UploadFile) -> bool:
    """Validate if uploaded file is a valid CSV"""
    if not file.filename.lower().endswith(('.csv', '.xlsx')):
        return False
    
    return True

def cleanup_session_files(session_id: str):
    """Clean up temporary files for a session"""
    directories = ["../data/uploads", "../data/outputs", "../data/plots"]
    
    for directory in directories:
        if os.path.exists(directory):
            for filename in os.listdir(directory):
                if filename.startswith(session_id):
                    file_path = os.path.join(directory, filename)
                    try:
                        os.remove(file_path)
                        print(f"Cleaned up: {file_path}")
                    except Exception as e:
                        print(f"Error cleaning up {file_path}: {e}")''',

    "frontend/app.py": '''import streamlit as st
import requests
import pandas as pd
import json
import time
import os
from io import BytesIO

st.set_page_config(
    page_title="DataSlide - CSV to PPT Converter",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="collapsed"
)

st.markdown("""
<style>
    .main > div {
        padding-top: 2rem;
    }
    
    .stApp > header {
        background-color: transparent;
    }
    
    .header-container {
        background: linear-gradient(90deg, #2C6AA0 0%, #1e4a72 100%);
        padding: 1rem 2rem;
        margin: -1rem -1rem 2rem -1rem;
        border-radius: 0;
    }
    
    .header-title {
        color: white;
        font-size: 2rem;
        font-weight: bold;
        margin: 0;
        display: flex;
        align-items: center;
        gap: 0.5rem;
    }
    
    .upload-container {
        border: 2px dashed #2C6AA0;
        border-radius: 10px;
        padding: 2rem;
        text-align: center;
        background-color: #f8f9fa;
        margin: 1rem 0;
    }
    
    .theme-card {
        border: 2px solid #e0e0e0;
        border-radius: 8px;
        padding: 1rem;
        text-align: center;
        cursor: pointer;
        transition: all 0.3s ease;
    }
    
    .theme-card:hover {
        border-color: #2C6AA0;
        box-shadow: 0 4px 8px rgba(44, 106, 160, 0.2);
    }
    
    .theme-card.selected {
        border-color: #2C6AA0;
        background-color: #e3f2fd;
    }
    
    .generate-button {
        background: linear-gradient(90deg, #2C6AA0 0%, #1e4a72 100%);
        color: white;
        border: none;
        padding: 0.75rem 2rem;
        border-radius: 8px;
        font-size: 1.1rem;
        font-weight: bold;
        cursor: pointer;
        width: 100%;
        margin: 1rem 0;
    }
</style>
""", unsafe_allow_html=True)

API_BASE_URL = "http://localhost:8000"

def main():
    st.markdown("""
    <div class="header-container">
        <div class="header-title">
            📊 DataSlide
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    if 'session_id' not in st.session_state:
        st.session_state.session_id = None
    if 'csv_data' not in st.session_state:
        st.session_state.csv_data = None
    if 'columns_info' not in st.session_state:
        st.session_state.columns_info = []
    if 'selected_columns' not in st.session_state:
        st.session_state.selected_columns = []
    if 'use_template' not in st.session_state:
        st.session_state.use_template = False
    if 'selected_theme' not in st.session_state:
        st.session_state.selected_theme = "Corporate Blue"
    
    col1, col2, col3 = st.columns([1, 2, 1])
    
    with col1:
        st.markdown("### 📁 Upload CSV File")
        csv_uploader()
        
    with col2:
        st.markdown("### 📋 Data Preview & Column Selection")
        data_preview()
        
    with col3:
        st.markdown("### 🎨 Template & Theme Selection")
        template_theme_selection()
    
    st.markdown("---")
    progress_section()
    
def csv_uploader():
    uploaded_file = st.file_uploader(
        "Choose a CSV file",
        type=['csv', 'xlsx'],
        help="Upload your CSV or Excel file to convert to PowerPoint"
    )
    
    if uploaded_file is not None:
        if st.session_state.session_id is None or uploaded_file.name != getattr(st.session_state, 'uploaded_filename', ''):
            with st.spinner("Uploading and analyzing file..."):
                upload_response = upload_csv_file(uploaded_file)
                if upload_response:
                    st.session_state.session_id = upload_response['session_id']
                    st.session_state.csv_data = upload_response
                    st.session_state.columns_info = upload_response['columns']
                    st.session_state.uploaded_filename = uploaded_file.name
                    st.success(f"✅ File uploaded successfully! ({upload_response['shape'][0]} rows, {upload_response['shape'][1]} columns)")
                    st.rerun()

def upload_csv_file(uploaded_file):
    try:
        files = {"file": (uploaded_file.name, uploaded_file.getvalue(), uploaded_file.type)}
        response = requests.post(f"{API_BASE_URL}/upload-csv", files=files)
        if response.status_code == 200:
            return response.json()
        else:
            st.error(f"Error uploading file: {response.text}")
            return None
    except Exception as e:
        st.error(f"Connection error: {e}")
        return None

def data_preview():
    if st.session_state.csv_data:
        st.markdown("**Data Preview:**")
        preview_df = pd.DataFrame(st.session_state.csv_data['preview'])
        st.dataframe(preview_df.head(), use_container_width=True)
        
        st.markdown("**Select Columns for Analysis:**")
        
        select_all = st.checkbox("Select All Columns")
        
        if select_all:
            st.session_state.selected_columns = [col['name'] for col in st.session_state.columns_info]
        
        selected_cols = []
        for col_info in st.session_state.columns_info:
            col_name = col_info['name']
            is_selected = st.checkbox(
                f"{col_name} ({col_info['data_type']})", 
                value=col_name in st.session_state.selected_columns or select_all,
                key=f"col_{col_name}"
            )
            if is_selected:
                selected_cols.append(col_name)
        
        st.session_state.selected_columns = selected_cols
        
        if selected_cols:
            st.info(f"📊 Selected {len(selected_cols)} columns for analysis")
    else:
        st.info("Upload a CSV file to see data preview and select columns")

def template_theme_selection():
    use_template = st.toggle("Use Template", value=st.session_state.use_template)
    st.session_state.use_template = use_template
    
    if use_template:
        st.markdown("**Upload PowerPoint Template:**")
        template_file = st.file_uploader(
            "Upload .pptx template",
            type=['pptx'],
            help="Upload your custom PowerPoint template"
        )
        
        if template_file and st.session_state.session_id:
            if st.button("Upload Template"):
                with st.spinner("Uploading template..."):
                    upload_response = upload_template(template_file)
                    if upload_response:
                        st.success("✅ Template uploaded successfully!")
    else:
        st.markdown("**Select Theme:**")
        themes = ["Corporate Blue", "Modern Green", "Creative Orange", "Minimal Gray"]
        
        selected_theme = st.selectbox(
            "Choose a theme",
            themes,
            index=themes.index(st.session_state.selected_theme)
        )
        st.session_state.selected_theme = selected_theme
        
        theme_colors = {
            "Corporate Blue": "#2C6AA0",
            "Modern Green": "#4CAF50", 
            "Creative Orange": "#FF9800",
            "Minimal Gray": "#607D8B"
        }
        
        st.markdown(f"""
        <div style="background-color: {theme_colors[selected_theme]}; color: white; padding: 1rem; border-radius: 8px; text-align: center; margin: 1rem 0;">
            <strong>{selected_theme}</strong> Theme Preview
        </div>
        """, unsafe_allow_html=True)

def upload_template(template_file):
    try:
        files = {"file": (template_file.name, template_file.getvalue(), template_file.type)}
        data = {"session_id": st.session_state.session_id}
        response = requests.post(f"{API_BASE_URL}/upload-template", files=files, data=data)
        if response.status_code == 200:
            return response.json()
        else:
            st.error(f"Error uploading template: {response.text}")
            return None
    except Exception as e:
        st.error(f"Connection error: {e}")
        return None

def progress_section():
    steps = ["1. Upload", "2. Settings", "3. Convert"]
    current_step = 0
    
    if st.session_state.session_id:
        current_step = 1
        if st.session_state.selected_columns or st.session_state.use_template:
            current_step = 2
    
    progress_cols = st.columns(3)
    for i, step in enumerate(steps):
        with progress_cols[i]:
            if i <= current_step:
                st.markdown(f"✅ **{step}**")
            else:
                st.markdown(f"⭕ {step}")
    
    st.markdown("---")
    
    can_generate = (
        st.session_state.session_id and 
        (st.session_state.selected_columns or len(st.session_state.columns_info) > 0)
    )
    
    if can_generate:
        if st.button("🚀 Generate Presentation", type="primary", use_container_width=True):
            generate_presentation()
    else:
        st.button("🚀 Generate Presentation", disabled=True, use_container_width=True)
        if not st.session_state.session_id:
            st.info("Please upload a CSV file first")

def generate_presentation():
    with st.spinner("🔄 Generating presentation... This may take a few minutes."):
        
        data = {
            "session_id": st.session_state.session_id,
            "selected_columns": json.dumps(st.session_state.selected_columns),
            "use_template": st.session_state.use_template,
            "theme": st.session_state.selected_theme
        }
        
        try:
            response = requests.post(f"{API_BASE_URL}/generate-presentation", data=data)
            
            if response.status_code == 200:
                result = response.json()
                st.success("✅ Presentation generated successfully!")
                
                if st.button("📥 Download Presentation", type="primary"):
                    download_presentation()
                    
            else:
                st.error(f"Error generating presentation: {response.text}")
                
        except Exception as e:
            st.error(f"Connection error: {e}")

def download_presentation():
    try:
        response = requests.get(f"{API_BASE_URL}/download/{st.session_state.session_id}")
        if response.status_code == 200:
            st.download_button(
                label="📥 Download PowerPoint",
                data=response.content,
                file_name=f"DataSlide_Presentation_{st.session_state.session_id[:8]}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        else:
            st.error("Error downloading presentation")
    except Exception as e:
        st.error(f"Download error: {e}")

if __name__ == "__main__":
    main()''',

    "run.py": '''#!/usr/bin/env python3
"""
DataSlide Project Runner
Automatically starts both backend and frontend servers
"""

import subprocess
import sys
import time
import os
import platform

def install_requirements():
    """Install required packages"""
    print("📦 Installing requirements...")
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "-r", "requirements.txt"])
        print("✅ Requirements installed successfully!")
    except subprocess.CalledProcessError as e:
        print(f"❌ Error installing requirements: {e}")
        return False
    return True

def start_backend():
    """Start FastAPI backend"""
    print("🚀 Starting FastAPI backend...")
    os.chdir("backend")
    
    if platform.system() == "Windows":
        return subprocess.Popen([sys.executable, "main.py"], shell=True)
    else:
        return subprocess.Popen([sys.executable, "main.py"])

def start_frontend():
    """Start Streamlit frontend"""
    print("🎨 Starting Streamlit frontend...")
    os.chdir("../frontend")
    
    if platform.system() == "Windows":
        return subprocess.Popen([sys.executable, "-m", "streamlit", "run", "app.py"], shell=True)
    else:
        return subprocess.Popen([sys.executable, "-m", "streamlit", "run", "app.py"])

def main():
    print("🌟 DataSlide - CSV to PPT Converter")
    print("=====================================")
    
    if not install_requirements():
        return
    
    try:
        backend_process = start_backend()
        time.sleep(5)
        
        frontend_process = start_frontend()
        
        print("\\n✅ Both servers are starting...")
        print("📱 Frontend: http://localhost:8501")
        print("⚡ Backend: http://localhost:8000")
        print("\\n🔄 Waiting for servers to initialize...")
        print("📝 Check the terminal output above for any errors")
        print("\\n❌ Press Ctrl+C to stop both servers")
        
        try:
            backend_process.wait()
            frontend_process.wait()
        except KeyboardInterrupt:
            print("\\n🛑 Stopping servers...")
            backend_process.terminate()
            frontend_process.terminate()
            print("✅ Servers stopped successfully!")
            
    except Exception as e:
        print(f"❌ Error starting servers: {e}")

if __name__ == "__main__":
    main()''',

    "sample_data.csv": '''Name,Position,Department,Salary,Years_Experience,Performance_Score
John Smith,Manager,Sales,85000,8,4.2
Jane Doe,Analyst,Finance,65000,3,4.5
Michael Brown,Engineer,Engineering,95000,5,4.8
Emily Davis,Designer,Marketing,72000,4,4.1
Anna Wilson,Consultant,Consulting,105000,7,4.6
David Johnson,Developer,Engineering,88000,6,4.4
Sarah Miller,Coordinator,Marketing,58000,2,3.9
James Taylor,Specialist,Finance,71000,4,4.3
Lisa Anderson,Manager,Sales,92000,9,4.7
Robert Wilson,Analyst,Engineering,78000,3,4.0
Maria Garcia,Designer,Marketing,69000,3,4.2
Christopher Lee,Developer,Engineering,91000,5,4.5
Jennifer Davis,Coordinator,Sales,62000,2,3.8
Matthew Brown,Specialist,Finance,76000,5,4.1
Amanda Johnson,Manager,Marketing,89000,7,4.4
Kevin Anderson,Engineer,Engineering,97000,6,4.6
Michelle Wilson,Analyst,Consulting,82000,4,4.3
Daniel Martinez,Developer,Engineering,93000,7,4.8
Ashley Taylor,Designer,Marketing,71000,4,4.0
Ryan Thomas,Specialist,Finance,74000,3,4.2''',

    "backend/__init__.py": "",
    "backend/services/__init__.py": "",
    "backend/models/__init__.py": "",
    "backend/utils/__init__.py": "",
    "data/uploads/.gitkeep": "",
    "data/outputs/.gitkeep": "",
    "data/plots/.gitkeep": "",
    "templates/.gitkeep": ""
}

# Combine all files
all_project_files = {
    "requirements.txt": """streamlit==1.29.0
fastapi==0.104.1
uvicorn==0.24.0
pandas==2.1.3
numpy==1.25.2
matplotlib==3.8.2
seaborn==0.13.0
python-pptx==0.6.23
python-multipart==0.0.6
requests==2.31.0
transformers==4.36.0
torch==2.1.1
accelerate==0.25.0
scipy==1.11.4
scikit-learn==1.3.2
plotly==5.17.0
openpyxl==3.1.2
Pillow==10.1.0
aiofiles==23.2.1""",
    **all_files_continued
}

print(f"Total project files to include: {len(all_project_files)}")

# Create the zip file with all project files
zip_filename = "DataSlide_Complete_Project.zip"

with zipfile.ZipFile(zip_filename, 'w', zipfile.ZIP_DEFLATED) as zipf:
    for file_path, content in all_project_files.items():
        zipf.writestr(file_path, content)
        print(f"✅ Added: {file_path}")

# Get final zip file info
zip_size = os.path.getsize(zip_filename)
zip_size_kb = zip_size / 1024

print(f"\n🎉 COMPLETE ZIP FILE CREATED!")
print(f"📦 Filename: {zip_filename}")
print(f"📊 Size: {zip_size_kb:.1f} KB")
print(f"📁 Contains: {len(all_project_files)} files")
print(f"✅ Ready for download!")